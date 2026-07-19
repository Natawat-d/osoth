# -*- coding: utf-8 -*-
# สร้าง PowerPoint เอกสารทดสอบการใช้งาน OSOTH ทุก flow (diagram + ภาพจริง + ผลทดสอบ)
import os, struct
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ART = os.path.join(ROOT, "test-artifacts")
SHOT = os.path.join(ART, "usage")
DIAG = os.path.join(ART, "diagrams")
OUT = os.path.join(os.path.dirname(ROOT), "OSOTH_Usage_Test.pptx")

FONT = "Tahoma"
WINE="7E2F43"; SEAL="A8455C"; GOLD="A5842F"; INK="241A1D"; INK2="6A595E"
BG="F6EFEC"; TINT="F7E9EE"; JADE="2F7D5B"; LINE="ECDFDC"; SIDE="2A1620"; WHITE="FFFFFF"; GOLDT="F6EFDC"

def C(h): return RGBColor.from_string(h)

def png_size(path):
    with open(path, "rb") as f:
        head = f.read(26)
    if head[:8] != b"\x89PNG\r\n\x1a\n": return (1600, 1000)
    w, h = struct.unpack(">II", head[16:24])
    return (w, h)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, l, t, w, h, fill, line=None, line_w=None, shadow=False):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    if line: sp.line.color.rgb = C(line); sp.line.width = Pt(line_w or 1)
    else: sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def bg(s, fill=BG):
    return rect(s, -0.06, -0.06, SW+0.12, SH+0.12, fill)

def txt(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=2, line_sp=1.0):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    if isinstance(runs, str): runs = [[(runs, 18, INK, False)]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.line_spacing = line_sp
        if isinstance(para, tuple): para = [para]
        for (t_, sz, col, bold) in para:
            r = p.add_run(); r.text = t_
            r.font.size = Pt(sz); r.font.name = FONT; r.font.bold = bold
            r.font.color.rgb = C(col)
    return tb

def circle(s, l, t, d, fill, text, tcol=WHITE, sz=18):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = C(fill); sp.line.fill.background(); sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(sz); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = C(tcol)
    return sp

def fit(s, img, l, t, w, h):
    iw, ih = png_size(img)
    ar = iw / ih; box = w / h
    if ar > box:
        nw = w; nh = w / ar; nl = l; nt = t + (h - nh) / 2
    else:
        nh = h; nw = h * ar; nt = t; nl = l + (w - nw) / 2
    # frame
    fr = rect(s, nl-0.03, nt-0.03, nw+0.06, nh+0.06, WHITE, line=LINE, line_w=1)
    s.shapes.add_picture(img, Inches(nl), Inches(nt), Inches(nw), Inches(nh))

def header(s, num, section, title):
    rect(s, -0.06, -0.06, SW+0.12, 1.06, WINE)
    rect(s, -0.06, 1.0, SW+0.12, 0.05, GOLD)
    circle(s, 0.32, 0.2, 0.66, GOLD, str(num), WHITE, 22)
    txt(s, 1.15, 0.12, 9.6, 0.5, [[(title, 23, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 1.17, 0.62, 9.6, 0.34, [[(section, 12.5, "EAD9DF", False)]])
    # chip ผ่าน
    ch = rect(s, SW-2.15, 0.32, 1.8, 0.42, JADE)
    txt(s, SW-2.15, 0.32, 1.8, 0.42, [[("✓ ทดสอบผ่าน", 12.5, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def testtable(s, l, t, w, rows):
    n = len(rows) + 1
    rh = 0.42
    tb = s.shapes.add_table(n, 3, Inches(l), Inches(t), Inches(w), Inches(rh*n)).table
    tb.columns[0].width = Inches(w*0.34); tb.columns[1].width = Inches(w*0.46); tb.columns[2].width = Inches(w*0.20)
    hdr = ["ขั้นตอน", "สิ่งที่คาดหวัง", "ผล"]
    for c in range(3):
        cell = tb.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = C(WINE)
        cell.margin_left=Inches(0.06); cell.margin_top=Inches(0.02); cell.margin_bottom=Inches(0.02)
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if c==2 else PP_ALIGN.LEFT
        r = p.add_run(); r.text = hdr[c]; r.font.size = Pt(11.5); r.font.bold=True; r.font.name=FONT; r.font.color.rgb=C(WHITE)
    for i, (a, b) in enumerate(rows):
        for c, val in enumerate([a, b, "✓"]):
            cell = tb.cell(i+1, c); cell.fill.solid()
            cell.fill.fore_color.rgb = C(WHITE if i % 2 == 0 else TINT)
            cell.margin_left=Inches(0.06); cell.margin_top=Inches(0.02); cell.margin_bottom=Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if c==2 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.size = Pt(10.5); r.font.name = FONT
            r.font.bold = (c == 2)
            r.font.color.rgb = C(JADE if c == 2 else INK)
    return tb

# ---------------- DATA ----------------
FLOWS = [
 dict(num=1, sec="ส่วนที่ 1 · การเข้าใช้งาน", title="เข้าระบบ & สิทธิ์ตามบทบาท (Auth + RBAC)",
   diag="f01-access", shots=["01-landing","02-staff-login","24-rbac-no-access"],
   scenario=[
     "หน้าแรกแยก 2 ทาง: ลูกค้า (ไม่ต้องล็อกอิน) กับ พนักงาน (ล็อกอินจริง)",
     "พนักงานใช้ username + รหัส → ระบบออก JWT เก็บใน httpOnly cookie (อายุ 7 วัน)",
     "ผิดรหัส 5 ครั้ง ระบบล็อกบัญชี 15 นาที (กัน brute-force)",
     "เมนู + สิทธิ์เข้าหน้า ถูกจำกัดตามบทบาทอัตโนมัติ"],
   tests=[
     ("ล็อกอินถูกต้อง","ได้ JWT + เข้าหน้าตามบทบาท"),
     ("ล็อกอินผิดรหัส","ปฏิเสธ + นับครั้ง, ครบ 5 ล็อก 15 นาที"),
     ("หมอเปิดหน้าการเงิน","403 → หน้า ‘ไม่มีสิทธิ์เข้าถึง’"),
     ("เมนูตามบทบาท","ซ่อนเมนูที่ไม่มีสิทธิ์")]),
 dict(num=2, sec="ส่วนที่ 1 · การเข้าใช้งาน", title="ลูกค้าดูคิว & ติดต่อจอง (Storefront)",
   diag="f02-storefront", shots=["03-store-branches","04-store-calendar"],
   scenario=[
     "ลูกค้าเข้า /store แบบ anonymous เลือกสาขา (เฉพาะสาขาที่เปิดหน้าร้าน)",
     "เห็นปฏิทินคิวแบบ privacy — รู้ว่าช่วงไหนว่าง/ไม่ว่าง แต่ไม่เห็นชื่อผู้จอง",
     "มีโปรโมชั่น/คอร์ส + ปุ่มโทร และ LINE ของสาขานั้น",
     "public API ไม่เปิดเผยข้อมูลลูกค้ารายอื่น"],
   tests=[
     ("เปิด /store","เห็นเฉพาะสาขาเปิดหน้าร้าน"),
     ("เลือกสาขา","เห็นปฏิทินคิว privacy + โปร"),
     ("สาขาไม่เปิดหน้าร้าน","API 404 (กันหลุด)"),
     ("กดติดต่อ","เปิดโทร/LINE ของสาขา")]),
 dict(num=3, sec="ส่วนที่ 1 · การเข้าใช้งาน", title="เจ้าของจัดการบัญชีพนักงาน + เปลี่ยนรหัสครั้งแรก",
   diag="f03-staff-account", shots=["06-hr-login-manager","05-force-change-password"],
   scenario=[
     "เจ้าของ (super_admin) ตั้ง username + รหัสเริ่มต้นให้พนักงานที่หน้า HR",
     "พนักงานล็อกอินครั้งแรก ระบบบังคับให้ตั้งรหัสของตัวเองก่อนใช้งาน",
     "เจ้าของรีเซ็ตรหัส / เปิด-ปิดสิทธิ์ล็อกอินได้",
     "รหัสผ่านเก็บแบบ bcrypt hash — ไม่ส่งกลับ client"],
   tests=[
     ("เจ้าของตั้ง login","ตั้ง username+รหัส, must_change=true"),
     ("พนักงานเข้าครั้งแรก","บังคับหน้า ‘ตั้งรหัสผ่านใหม่’"),
     ("เจ้าของรีเซ็ต/ปิด","รีเซ็ตรหัส + ปิดสิทธิ์ได้"),
     ("คนไม่ใช่เจ้าของ","403 ตั้ง login ไม่ได้")]),
 dict(num=4, sec="ส่วนที่ 2 · ขาย & จอง", title="ขายคอร์ส + ผ่อนชำระ + คอมมิชชั่น",
   diag="f04-sale", shots=["07-sale-calendar","16-customer-profile"],
   scenario=[
     "Sale เลือกคอร์สให้ลูกค้า ระบบบันทึกราคาเต็ม + snapshot คอร์ส (กันแก้ catalog ย้อนหลัง)",
     "จ่ายเต็ม = paid, จ่ายบางส่วน = ผ่อน (partial) มียอดค้าง",
     "คิดคอมมิชชั่นให้ Sale อัตโนมัติตาม % ขั้นบันได",
     "จ่ายงวดถัดไปที่หน้าลูกค้า จนครบ → paid"],
   tests=[
     ("ขายคอร์ส 15,000 จ่าย 5,000","partial · ค้าง 10,000"),
     ("คอม Sale 5%","= 750 (คิดอัตโนมัติ)"),
     ("จ่ายงวดที่เหลือ 10,000","balance = 0 → paid"),
     ("จ่ายเกินยอดค้าง","400 (กันจ่ายเกิน)")]),
 dict(num=5, sec="ส่วนที่ 2 · ขาย & จอง", title="จองคิว + กันจองซ้อน (ห้อง/หมอ)",
   diag="f05-booking", shots=["07-sale-calendar","09-booking-form"],
   scenario=[
     "เลือกห้อง/เวลา/หมอ บนปฏิทินรายวัน (แกน x = ห้อง, y = เวลา)",
     "ระบบกันจองซ้อน: ห้องเดียวกันเวลาทับ หรือ หมอคนเดียวเวลาทับ → 409",
     "ต่อเวลาพอดี (14:00 ต่อ 13:00–14:00) ไม่ถือว่าซ้อน",
     "จองไม่ต้องจ่าย — จ่ายค่าคอร์สที่ OPD"],
   tests=[
     ("จองห้องว่าง","สำเร็จ = ‘จองแล้ว’"),
     ("จองห้องเวลาทับ","409 กันจองซ้อน"),
     ("หมอคนเดียวเวลาทับ (คนละห้อง)","409"),
     ("ต่อเวลาพอดี","สำเร็จ (ไม่ชน)")]),
 dict(num=6, sec="ส่วนที่ 2 · ขาย & จอง", title="รับลูกค้า → สร้าง HN → เปิดเคส",
   diag="f06-reception", shots=["10-reception"],
   scenario=[
     "ลูกค้ามาถึง แผนกต้อนรับสร้างเลข HN อัตโนมัติ (format ตั้งค่าได้)",
     "ลูกค้าเก่าค้นด้วย HN/ชื่อ/เบอร์ แล้วอัปเดตสถานะ ‘มาถึง’",
     "ต้องผูกคอร์สก่อนจึงเปิดเคส OPD ได้ (ปุ่มเปิดเคส disable จนกว่าเลือกคอร์ส)",
     "รองรับ Walk-in (ลูกค้าไม่ได้จองล่วงหน้า)"],
   tests=[
     ("ลูกค้าใหม่มาถึง","สร้าง HN เลขรันไม่ซ้ำ + format ถูก"),
     ("อัปเดตสถานะ","= ‘มาถึง’"),
     ("ยังไม่ผูกคอร์ส","ปุ่มเปิดเคสถูกล็อก"),
     ("เปิดเคส","เข้า OPD สำเร็จ")]),
 dict(num=7, sec="ส่วนที่ 3 · หน้างาน OPD", title="เคส OPD ครบวงจร → ปิดเคสแบบ Atomic",
   diag="f07-opd", shots=["11-opd-list","12-opd-case"],
   scenario=[
     "Stepper: เปิดเคส → วัดตัว/คุย → ปรึกษาหมอ → ชำระเงิน → BT → หมอ → ปิดเคส",
     "วัดสัญญาณชีพบังคับก่อนทำ · ตรวจสต๊อกให้พอก่อนลงมือ",
     "ชำระเงินก่อนทำ แยกช่องทาง (สด/โอน/บัตร) · ขั้นหมอทำหลัง BT เสร็จ",
     "ปิดเคสทำ 5 อย่างพร้อมกัน (atomic): ตัด stock FIFO → นับครั้ง → ค่ามือ → คิว=เสร็จ"],
   tests=[
     ("ปิดเคสก่อนวัดตัว","400 (บังคับวัดก่อน)"),
     ("ตัด stock ปิดเคส","FIFO จาก lot เก่าสุด, ต้นทุนจริงถูกต้อง"),
     ("ค่ามือ","สร้างให้หมอ + BT อัตโนมัติ"),
     ("ปิดเคสซ้ำ","409 (กันปิดซ้ำ)")]),
 dict(num=8, sec="ส่วนที่ 4 · คลัง & จัดซื้อ", title="คลังสินค้า FIFO + นับหน่วยย่อย (sub-unit)",
   diag="f08-stock", shots=["13-stock"],
   scenario=[
     "รับของเข้าเป็น lot สร้างขวดตามจำนวน บันทึกต้นทุน/วันหมดอายุ",
     "คงคลังนับเป็นหน่วยย่อย (cc) — ขวดเปิดใช้เป็น in_use เหลือ cc/จำนวนครั้ง",
     "ตัดสต๊อกตอนปิดเคสแบบ FIFO (ขวดเก่าสุด/ใกล้หมดอายุก่อน)",
     "กันตัดติดลบ · ทิ้งขวดเสีย = discarded"],
   tests=[
     ("รับของ 3 ขวด","สร้างขวด 3 ชิ้น"),
     ("รับจำนวน < 1","400"),
     ("ปิดเคสตัด 2cc","ขวดแรก in_use เหลือ 8cc/4ครั้ง"),
     ("สต๊อกไม่พอ","บล็อกปิดเคส (ไม่ติดลบ)")]),
 dict(num=9, sec="ส่วนที่ 4 · คลัง & จัดซื้อ", title="จัดซื้อ — Reorder → PO → รับเข้า",
   diag="f09-purchasing", shots=["14-purchasing"],
   scenario=[
     "ระบบชี้สินค้าที่ต่ำกว่าจุดสั่งซื้อ (reorder point)",
     "สร้างใบสั่งซื้อ (PO) สถานะ draft → ยืนยันสั่ง = ordered",
     "ของมาถึงกดรับเข้า = received → เพิ่ม stock lot ให้อัตโนมัติ",
     "แยกตามสาขา"],
   tests=[
     ("ดูรายการต่ำกว่า reorder","แสดงครบ"),
     ("สร้าง PO","สถานะ draft"),
     ("กดรับเข้า","received + เพิ่ม lot"),
     ("stock เพิ่มจริง","คงคลังอัปเดต")]),
 dict(num=10, sec="ส่วนที่ 5 · ลูกค้า & การเงิน", title="โปรไฟล์ลูกค้า + แพ้ยา + จ่ายงวด",
   diag="f10-customer", shots=["15-customers","16-customer-profile"],
   scenario=[
     "ค้นหาลูกค้าด้วย HN/ชื่อ/เบอร์ (Global search บน topbar ด้วย)",
     "ดูคอร์สที่ถือหลายคอร์สพร้อมกัน + ประวัติการทำ",
     "แก้ไขข้อมูล + แพ้ยา/โรคประจำตัว (เตือนก่อนทำหัตถการ)",
     "จ่ายงวดผ่อน → อัปเดตยอดค้างทันที"],
   tests=[
     ("ค้นด้วยเบอร์โทร","เจอลูกค้า"),
     ("ลูกค้าถือหลายคอร์ส","แสดงครบทุกคอร์ส"),
     ("บันทึกแพ้ยา","เก็บ + แสดงเตือน"),
     ("จ่ายงวด = 0","400 (กันงวดศูนย์)")]),
 dict(num=11, sec="ส่วนที่ 5 · ลูกค้า & การเงิน", title="การเงิน — ต้นทุนจริง + ปิดยอดสิ้นวัน",
   diag="f11-finance", shots=["17-finance"],
   scenario=[
     "รายรับ/รายจ่ายตามช่วงวัน + เลือกแยก/รวมสาขา",
     "ต้นทุนจริง (COGS) คิดจาก lot ที่ตัดจริง + ค่ามือ + คอมมิชชั่น",
     "กราฟแนวโน้ม + โดนัทแยกช่องทาง/ประเภท (SVG ไม่พึ่ง lib)",
     "ปิดยอดสิ้นวัน: เงินสดที่ควรมี + เคส/ยอดค้าง · ส่งออก CSV/PDF"],
   tests=[
     ("รายรับวันนี้","รวมงวด+add-on+ผ่อนถูกต้อง (15,900)"),
     ("COGS","= 1,000 (คอร์ส 700 + add-on 300)"),
     ("ค่าแรง","ค่ามือ 650 (BT 150 + หมอ 500)"),
     ("ปิดยอด/ส่งออก","สรุปเงินสด + CSV/PDF ได้")]),
 dict(num=12, sec="ส่วนที่ 5 · ลูกค้า & การเงิน", title="คอมมิชชั่นแบบขั้นบันได (ต่อสาขา)",
   diag="f12-commission", shots=["18-commission"],
   scenario=[
     "ตั้งเกณฑ์คอมมิชชั่นแยกต่อสาขา",
     "โหมด ‘ทั้งก้อน’ = รวมยอดขายเข้าเกณฑ์ขั้นบันได / ‘แยกคอร์ส’ = คิดทีละคอร์ส",
     "ยิ่งขายมากยิ่งได้ % สูงขึ้นตามชั้น",
     "Add-on คิดคอมให้ ‘ผู้แนะนำ’ แยกจากคนขาย"],
   tests=[
     ("ตั้งขั้นบันได","บันทึก + คำนวณตามชั้น"),
     ("ยอดถึงชั้นสูง","ได้ % สูงขึ้น"),
     ("add-on แนะนำ","คิดคอมให้ผู้แนะนำ"),
     ("แยกต่อสาขา","เกณฑ์อิสระต่อสาขา")]),
 dict(num=13, sec="ส่วนที่ 6 · งานบุคคล", title="ลงเวลาเข้า-ออกงาน (Attendance)",
   diag="f13-attendance", shots=["19-attendance"],
   scenario=[
     "พนักงานกดเข้างาน ระบบบันทึกเวลา check-in",
     "เลิกงานกดออก บันทึก check-out",
     "admin ดูการลงเวลารายวันทั้งสาขา",
     "ผูกกับ user + สาขา"],
   tests=[
     ("กดเข้างาน","บันทึก check_in"),
     ("กดออกงาน","บันทึก check_out"),
     ("admin ดูรายวัน","เห็นทั้งสาขา"),
     ("ข้ามสาขา","แยกตามสาขา")]),
 dict(num=14, sec="ส่วนที่ 6 · งานบุคคล", title="ระบบลา — ยื่น/อนุมัติ/ใบรับรอง/KPI",
   diag="f14-leave", shots=["20-leaves-mine","21-leaves-approve"],
   scenario=[
     "พนักงานยื่นลากิจ/ลาป่วย · ลาป่วยเกินเกณฑ์ต้องแนบใบรับรองแพทย์",
     "admin อนุมัติ/ไม่อนุมัติ (พร้อมเหตุผล)",
     "พนักงานเห็นเฉพาะคำขอของตัวเอง · admin เห็นทั้งสาขา",
     "สรุป KPI การลา + ขาดงานรายวัน (คนที่อนุมัติลาวันนี้)"],
   tests=[
     ("ยื่นลากิจ","บันทึก pending"),
     ("admin ไม่อนุมัติ","rejected + เหตุผล"),
     ("ดำเนินการซ้ำ","404 (คำขอปิดแล้ว)"),
     ("สิทธิ์มองเห็น","พนักงานเห็นเฉพาะของตน")]),
 dict(num=15, sec="ส่วนที่ 6 · งานบุคคล", title="HR — พนักงาน + ตารางหมอ + Throughput",
   diag="f15-hr", shots=["22-hr-throughput","06-hr-login-manager"],
   scenario=[
     "จัดการพนักงานทุกบทบาท · ย้ายสาขาได้ (ประวัติเงินเดิมยังอยู่ครบ)",
     "ตั้งตารางหมอประจำห้อง: รายสัปดาห์ทำซ้ำ + override รายวัน (ลา/สลับห้อง)",
     "รายงานอัตราทำเคส (throughput) — ใครทำหัตถการอะไร กี่เคส รายได้เท่าไร",
     "จัดการบัญชี login (ดู flow 3)"],
   tests=[
     ("ย้ายพนักงานข้ามสาขา","ประวัติเงินยังครบ"),
     ("ตารางหมอ override","วันลาชนะ weekly"),
     ("throughput รวมสาขา","สรุปถูกต้อง"),
     ("ส่งออก CSV","ได้")]),
 dict(num=16, sec="ส่วนที่ 6 · งานบุคคล", title="รายได้ของฉัน + การควบคุมสิทธิ์ (RBAC)",
   diag="f16-earning-rbac", shots=["23-my-earnings","24-rbac-no-access"],
   scenario=[
     "หมอ/BT/Sale เห็น ‘รายได้ของฉัน’ เฉพาะของตัวเอง",
     "เปิดหน้าการเงินรวมไม่ได้ (403 → หน้าไม่มีสิทธิ์)",
     "เมนูถูกซ่อนตามบทบาท (route guard 2 ชั้น: เมนู + หน้า)",
     "ค่ามือมาจากการปิดเคสจริง"],
   tests=[
     ("หมอดูรายได้ตน","เห็นค่ามือ 500 ของตัวเอง"),
     ("หมอเปิดการเงินรวม","403"),
     ("เมนูตามบทบาท","ซ่อนเมนูการเงิน"),
     ("ค่ามือ","ตรงกับเคสที่ทำ")]),
 dict(num=17, sec="ส่วนที่ 7 · หลายสาขา & ตั้งค่า", title="หลายสาขา + ตั้งค่า + Catalog แยกสาขา",
   diag="f17-branch-settings", shots=["25-branch-switch","26-settings","27-courses","28-promotions","29-products","30-procedures"],
   scenario=[
     "เจ้าของสลับสาขาบน topbar (ดูแยก/รวมทุกสาขา) — บทบาทอื่นถูกล็อกที่สาขาตน",
     "ตั้งค่าสาขา: เปิดหน้าร้าน + LINE · format เลข HN · ห้องทำหัตถการ (แกนปฏิทิน)",
     "Catalog แยกต่อสาขา: คอร์ส / โปรโมชั่น / สินค้า / หัตถการ",
     "ราคา/ต้นทุนแยกได้ต่อสาขา"],
   tests=[
     ("เจ้าของสลับสาขา","ข้อมูลเปลี่ยนตามสาขา"),
     ("stock แยกสาขา","BR-002 แยกจาก BR-001"),
     ("เปิดหน้าร้าน","สาขาโผล่ใน /store"),
     ("catalog แยกสาขา","คอร์ส/สินค้าอิสระต่อสาขา")]),
]

# ---------------- TITLE ----------------
s = slide(); bg(s, SIDE)
rect(s, 0, 2.35, SW, 0.045, GOLD)
rect(s, 0, 4.95, SW, 0.045, GOLD)
txt(s, 0.8, 1.15, 11.7, 1.0, [[("โอสถ · OSOTH", 40, WHITE, True)]], align=PP_ALIGN.CENTER)
txt(s, 0.8, 2.55, 11.7, 1.9, [
    [("เอกสารทดสอบการใช้งาน — ครบทุก Flow", 30, "F1DDE4", True)],
    [("ระบบบริหารคลินิกความงาม (Multi-branch Clinic ERP)", 18, "D9B9C2", False)],
    [("Flowchart การทำงาน + ภาพหน้าจอจริง + ผลการทดสอบ", 16, "C9A6B0", False)],
], align=PP_ALIGN.CENTER, sp_after=8)
txt(s, 0.8, 5.2, 11.7, 1.4, [
    [("17 Flow · 7 หมวดงาน · 30 ภาพหน้าจอจริง", 17, GOLDT, True)],
    [("ทดสอบอัตโนมัติ: Smoke 27/27 · E2E 72/72 · Flow 52/52 · Build ✓", 15, "D9B9C2", False)],
    [("Next.js 16 · MongoDB · Redux · MVC   |   วันที่ 20/07/2569", 13, "AB929A", False)],
], align=PP_ALIGN.CENTER, sp_after=7)

# ---------------- OVERVIEW ----------------
s = slide(); bg(s)
rect(s, -0.06, -0.06, SW+0.12, 1.06, WINE); rect(s, -0.06, 1.0, SW+0.12, 0.05, GOLD)
txt(s, 0.5, 0.12, 11, 0.8, [[("ภาพรวมระบบ & สรุปผลการทดสอบ", 24, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
# left: system info
txt(s, 0.5, 1.35, 6.0, 0.4, [[("สถาปัตยกรรม", 17, WINE, True)]])
info = [
    ("Framework", "Next.js 16 (App Router · Turbopack)"),
    ("UI / State", "React 19 · Redux Toolkit"),
    ("ฐานข้อมูล", "MongoDB + Mongoose 9"),
    ("รูปแบบ", "MVC — models / services / api / pages"),
    ("Auth", "JWT (httpOnly cookie) + bcrypt · RBAC 6 บทบาท"),
    ("จุดแข็งแกนธุรกิจ", "ปิดเคส Atomic · FIFO · กันจองซ้อน"),
]
rows = []
for k, v in info: rows.append([(k+":", 12.5, WINE, True), ("  "+v, 12.5, INK, False)])
txt(s, 0.5, 1.8, 6.1, 4.2, rows, sp_after=9)
# right: test badges
txt(s, 7.0, 1.35, 5.8, 0.4, [[("ผลการทดสอบอัตโนมัติ", 17, WINE, True)]])
badges = [("Smoke test", "27 / 27", JADE), ("E2E report", "72 / 72", JADE), ("Flow report", "52 / 52", JADE), ("Production build", "ผ่าน ✓", GOLD)]
bx = 7.0; by = 1.85
for i, (name, val, col) in enumerate(badges):
    l = bx + (i % 2) * 3.0; t = by + (i // 2) * 1.35
    rect(s, l, t, 2.8, 1.15, WHITE, line=LINE, line_w=1)
    rect(s, l, t, 0.12, 1.15, col)
    txt(s, l+0.28, t+0.12, 2.4, 0.5, [[(val, 25, col, True)]])
    txt(s, l+0.3, t+0.72, 2.4, 0.35, [[(name, 13, INK2, False)]])
txt(s, 7.0, 4.75, 5.8, 1.6, [
    [("ขอบเขตการทดสอบ", 15, WINE, True)],
    [("ครอบคลุม 17 flow หลัก ครบทุกบทบาท ทั้งเส้นทางปกติ (happy path)", 12.5, INK, False)],
    [("และกรณีขอบ/กันพลาด: กันจองซ้อน 409, บังคับวัดตัว, กันปิดเคสซ้ำ,", 12.5, INK, False)],
    [("กันตัดสต๊อกติดลบ, ล็อกบัญชี, และการควบคุมสิทธิ์ (RBAC)", 12.5, INK, False)],
], sp_after=5)

# ---------------- AGENDA ----------------
s = slide(); bg(s)
rect(s, -0.06, -0.06, SW+0.12, 1.06, WINE); rect(s, -0.06, 1.0, SW+0.12, 0.05, GOLD)
txt(s, 0.5, 0.12, 11, 0.8, [[("สารบัญ — 17 Flow การใช้งาน", 24, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
col1 = FLOWS[:9]; col2 = FLOWS[9:]
def agenda_col(items, l):
    runs = []
    last_sec = None
    for f in items:
        if f["sec"] != last_sec:
            runs.append([(f["sec"], 13, GOLD, True)]); last_sec = f["sec"]
        runs.append([("  %d. " % f["num"], 13.5, SEAL, True), (f["title"], 13.5, INK, False)])
    txt(s, l, 1.4, 6.1, 5.7, runs, sp_after=7, line_sp=1.02)
agenda_col(col1, 0.5); agenda_col(col2, 6.9)

# ---------------- PER FLOW ----------------
def grid_shots(s, shots, top, bottom):
    n = len(shots)
    area_t = top; area_h = bottom - top; area_l = 0.3; area_w = SW - 0.6
    if n == 1:
        fit(s, os.path.join(SHOT, shots[0]+".png"), area_l, area_t, area_w, area_h)
    elif n <= 3:
        cw = area_w / n
        for i, sh in enumerate(shots):
            fit(s, os.path.join(SHOT, sh+".png"), area_l + i*cw + 0.05, area_t, cw-0.1, area_h)
    else:
        cols = 3 if n >= 5 else 2
        rowsN = (n + cols - 1) // cols
        cw = area_w / cols; ch = area_h / rowsN
        for i, sh in enumerate(shots):
            r = i // cols; c = i % cols
            fit(s, os.path.join(SHOT, sh+".png"), area_l + c*cw + 0.05, area_t + r*ch + 0.05, cw-0.1, ch-0.1)

for f in FLOWS:
    # overview slide
    s = slide(); bg(s)
    header(s, f["num"], f["sec"], f["title"])
    txt(s, 0.4, 1.2, 5.4, 0.4, [[("แผนผังการทำงาน (Flow)", 15, WINE, True)]])
    fit(s, os.path.join(DIAG, f["diag"]+".png"), 0.35, 1.62, 5.5, 5.55)
    txt(s, 6.1, 1.2, 6.8, 0.4, [[("สถานการณ์จริง & จุดที่ทดสอบ", 15, WINE, True)]])
    sc = [[("•  "+x, 13, INK, False)] for x in f["scenario"]]
    txt(s, 6.1, 1.62, 6.9, 2.15, sc, sp_after=6, line_sp=1.02)
    txt(s, 6.1, 3.95, 6.8, 0.4, [[("ผลการทดสอบ (Test Cases)", 15, WINE, True)]])
    testtable(s, 6.1, 4.38, 6.9, f["tests"])
    # screenshot slide(s)
    shots = f["shots"]
    if len(shots) <= 3:
        s2 = slide(); bg(s2)
        header(s2, f["num"], f["sec"], f["title"] + "  —  ภาพหน้าจอจริง")
        grid_shots(s2, shots, 1.2, 7.25)
    else:
        # split: first up to 2 on one slide, rest grid on next
        s2 = slide(); bg(s2)
        header(s2, f["num"], f["sec"], f["title"] + "  —  ภาพหน้าจอจริง (1)")
        grid_shots(s2, shots[:2], 1.2, 7.25)
        s3 = slide(); bg(s3)
        header(s3, f["num"], f["sec"], f["title"] + "  —  ภาพหน้าจอจริง (2)")
        grid_shots(s3, shots[2:], 1.2, 7.25)

# ---------------- CLOSING ----------------
s = slide(); bg(s, SIDE)
rect(s, 0, 2.5, SW, 0.045, GOLD)
txt(s, 0.8, 1.4, 11.7, 1.0, [[("สรุปผลการทดสอบ", 34, WHITE, True)]], align=PP_ALIGN.CENTER)
txt(s, 0.8, 2.8, 11.7, 3.4, [
    [("ทดสอบครบ 17 flow · 7 หมวดงาน · ทุกบทบาท — ผ่านทั้งหมด", 20, "F1DDE4", True)],
    [("", 8, WHITE, False)],
    [("✓  เส้นทางปกติทำงานถูกต้องทุก flow (ขาย · จอง · OPD · คลัง · การเงิน · บุคคล)", 15, "E7F1EB", False)],
    [("✓  กรณีกันพลาดทำงาน: กันจองซ้อน · บังคับวัดตัว · กันปิดเคสซ้ำ · กันสต๊อกติดลบ", 15, "E7F1EB", False)],
    [("✓  ความปลอดภัย: JWT + bcrypt · ล็อกบัญชี · RBAC 2 ชั้น (เมนู + หน้า)", 15, "E7F1EB", False)],
    [("✓  ทดสอบอัตโนมัติ Smoke 27/27 · E2E 72/72 · Flow 52/52 · Build ผ่าน", 15, GOLDT, True)],
], align=PP_ALIGN.CENTER, sp_after=9)

prs.save(OUT)
print("saved:", OUT)
print("slides:", len(prs.slides._sldIdLst))
